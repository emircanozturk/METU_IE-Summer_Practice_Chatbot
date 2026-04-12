import streamlit as st
from openai import OpenAI
import os
from pypdf import PdfReader
from pptx import Presentation
from docx import Document

# --- PAGE CONFIGURATION & CUSTOM CSS ---
st.set_page_config(page_title="METU-IE Summer Practice Bot", page_icon="🔴", layout="centered")

st.markdown("""
    <style>
    .block-container { 
        padding-top: 2rem; 
        max-width: 95% !important; 
    }
    
    [data-testid="stSidebar"] {
        min-width: 330px !important;
        max-width: 330px !important;
    }
    
    [data-testid="stSidebar"] .stImage {
        margin-top: -5px;
        margin-bottom: -5px;
    }
    [data-testid="stSidebar"] .stRadio {
        margin-top: -15px;
        margin-bottom: -10px;
    }
    
    div[role="radiogroup"] {
        gap: 10px !important;
    }

    .metu-topbar {
        background-color: #333333;
        color: #dddddd;
        padding: 6px 15px;
        font-size: 12px;
        border-radius: 6px 6px 0 0;
        text-align: center;
        font-family: Arial, sans-serif;
    }
    
    .metu-navbar {
        background-color: #cb2c30; 
        color: white;
        padding: 12px 0;
        text-align: center;
        font-weight: 700;
        font-size: 22px;
        margin-bottom: 30px;
        border-radius: 0 0 6px 6px;
        box-shadow: 0 4px 6px rgba(0,0,0,0.1);
        letter-spacing: 0.5px;
    }
    
    .header-text {
        color: #2c3e50;
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        margin-top: 15px;
        font-size: 28px;
        font-weight: 600;
        line-height: 1.2;
        text-align: center;
    }

    .stButton>button {
        width: 100%;
        text-align: left;
        background-color: transparent;
        border: 1px solid #ccc;
        color: #2c3e50;
        transition: all 0.2s;
    }
    .stButton>button:hover {
        border-color: #cb2c30;
        color: #cb2c30;
    }
    
    .sidebar-link {
        display: block;
        padding: 10px 15px;
        margin-bottom: 10px;
        background-color: #f8f9fa;
        border-left: 5px solid #cb2c30;
        color: #333 !important;
        text-decoration: none;
        border-radius: 4px;
        font-weight: 500;
        transition: all 0.2s ease-in-out;
        box-shadow: 0 1px 3px rgba(0,0,0,0.05);
    }
    .sidebar-link:hover {
        background-color: #e9ecef;
        transform: translateX(3px);
    }
    
    /* Starter Question Buttons styling */
    .starter-btn>button {
        background-color: #f8f9fa;
        text-align: center;
        border-radius: 8px;
        color: #cb2c30;
        border: 1px solid #cb2c30;
        margin-bottom: 5px;
        font-size: 14px;
        padding: 8px;
    }
    .starter-btn>button:hover {
        background-color: #cb2c30;
        color: white;
    }
    </style>
""", unsafe_allow_html=True)

# --- BILINGUAL DICTIONARY ---
lang_dict = {
    "ENG": {
        "topbar": "IE | METU",
        "title": "Industrial Engineering<br>Summer Practice",
        "subtitle": "Intelligent Chatbot Assistant for IE300/IE400",
        "faq_header": "❓ Quick Questions",
        "links_header": "🔗 Quick Links",
        "contact_hdr": "📞 Contact Information",
        "link_curr": "IE Curriculum",
        "link_sp": "IE SP Webpage",
        "link_main": "IE Main Page",
        "cat1": "General & Paperwork",
        "cat2": "IE 300 & IE 400 Specifics",
        "q1_lbl": "How to arrange an organization?",
        "q1_prmpt": "How can I arrange a summer practice organization?",
        "q2_lbl": "What is the required paperwork?",
        "q2_prmpt": "What is the paperwork needed to get officially registered to summer practice?",
        "q3_lbl": "Which departments should I visit?",
        "q3_prmpt": "Which departments of the organization should I visit during my summer practice?",
        "q4_lbl": "Are banks suitable for IE 400?",
        "q4_prmpt": "What sort of branches or divisions in a bank are suitable for an IE 400 practice organization?",
        "q5_lbl": "Manufacturing rules in service firms?",
        "q5_prmpt": "How should a student answer questions in manufacturing oriented areas if the practice is done in a service organization?",
        "q6_lbl": "Sufficient problem definition (IE400)?",
        "q6_prmpt": "What constitutes a 'sufficient' IE problem definition in IE 400 reports?",
        "q7_lbl": "Can I do a full-time project instead?",
        "q7_prmpt": "If I am planning to work on a project full-time at the summer practice organization, is it possible for me to present the project work?",
        "q8_lbl": "Rules for multiple students at one firm?",
        "q8_prmpt": "What sort of reporting is acceptable if several students practice in the same organization?",
        "placeholder": "E.g., What are the requirements for IE 300?",
        "thinking": "Consulting official guidelines...",
        "welcome_msg": "👋 Welcome! I am ready to help. Try asking one of these:",
        "rec_1": "IE 300 Prerequisites",
        "rec_1_prmpt": "What are the exact prerequisites for IE 300?",
        "rec_2": "IE 400 Company Types",
        "rec_2_prmpt": "What types of companies are accepted for IE 400?",
        "rec_3": "Required Paperwork",
        "rec_3_prmpt": "What is the paperwork needed to officially register?",
        "rec_4": "Report Deadlines",
        "rec_4_prmpt": "When is the deadline to submit the summer practice report?",
        "rec_5": "Finding a Company",
        "rec_5_prmpt": "How can I arrange a summer practice organization?",
        "rec_6": "Full-time Project",
        "rec_6_prmpt": "Can I do a full-time project instead of a regular practice?"
    },
    "TR": {
        "topbar": "EM | ODTÜ",
        "title": "Endüstri Mühendisliği<br>Yaz Stajı",
        "subtitle": "IE300/IE400 için Akıllı Sohbet Robotu Asistanı",
        "faq_header": "❓ Hızlı Sorular",
        "links_header": "🔗 Hızlı Bağlantılar",
        "contact_hdr": "📞 İletişim Bilgileri",
        "link_curr": "EM Müfredatı",
        "link_sp": "EM Yaz Stajı Sayfası",
        "link_main": "EM Ana Sayfası",
        "cat1": "Genel & Evrak İşleri",
        "cat2": "IE 300 & IE 400 Detayları",
        "q1_lbl": "Staj yeri nasıl ayarlanır?",
        "q1_prmpt": "Yaz stajı organizasyonunu nasıl ayarlayabilirim?",
        "q2_lbl": "Gerekli evraklar nelerdir?",
        "q2_prmpt": "Yaz stajına resmi olarak kayıt olmak için gerekli evraklar nelerdir?",
        "q3_lbl": "Hangi departmanları ziyaret etmeliyim?",
        "q3_prmpt": "Yaz stajım sırasında organizasyonun hangi departmanlarını ziyaret etmeliyim?",
        "q4_lbl": "Bankalar IE 400 için uygun mu?",
        "q4_prmpt": "Bir bankadaki hangi tür şubeler veya bölümler IE 400 staj organizasyonu için uygundur?",
        "q5_lbl": "Hizmet firmalarında üretim soruları?",
        "q5_prmpt": "Staj hizmet organizasyonunda yapılıyorsa, öğrenci üretim odaklı alanlardaki soruları nasıl yanıtlamalıdır?",
        "q6_lbl": "Yeterli problem tanımı (IE400)?",
        "q6_prmpt": "IE 400 raporlarında 'yeterli' bir Endüstri Mühendisliği problem tanımı neleri içerir?",
        "q7_lbl": "Bunun yerine tam zamanlı proje yapabilir miyim?",
        "q7_prmpt": "Yaz stajı organizasyonunda tam zamanlı bir projede çalışmayı planlıyorsam, proje çalışmasını sunmam mümkün mü?",
        "q8_lbl": "Aynı firmadaki birden fazla öğrenci için kurallar?",
        "q8_prmpt": "Aynı organizasyonda birkaç öğrenci staj yaparsa ne tür bir raporlama kabul edilir?",
        "placeholder": "Örn., IE 300 için gereksinimler nelerdir?",
        "thinking": "Resmi yönergeler inceleniyor...",
        "welcome_msg": "👋 Hoş geldiniz! Yardım etmeye hazırım. Şunları sorabilirsiniz:",
        "rec_1": "IE 300 Ön Koşulları",
        "rec_1_prmpt": "IE 300 stajı için ön koşullar nelerdir?",
        "rec_2": "IE 400 Şirket Türleri",
        "rec_2_prmpt": "IE 400 stajı için hangi tür şirketler kabul ediliyor?",
        "rec_3": "Gerekli Evraklar",
        "rec_3_prmpt": "Resmi kayıt için gerekli evraklar nelerdir?",
        "rec_4": "Teslim Tarihleri",
        "rec_4_prmpt": "Yaz stajı raporunu teslim etmek için son tarih nedir?",
        "rec_5": "Şirket Bulma",
        "rec_5_prmpt": "Yaz stajı organizasyonunu nasıl ayarlayabilirim?",
        "rec_6": "Tam Zamanlı Proje",
        "rec_6_prmpt": "Normal staj yerine tam zamanlı bir proje yapabilir miyim?"
    }
}

# --- SIDEBAR & LANGUAGE SELECTION ---
with st.sidebar:
    selected_lang = st.radio("Language / Dil", ["ENG", "TR"], horizontal=True, label_visibility="collapsed")
    st.markdown("---")
    try:
        st.image("metu_logo.png", use_container_width=True)
    except:
        st.error("Logo missing. Add 'metu_logo.png' to folder.")
        
t = lang_dict[selected_lang]

# --- REUSABLE UI HEADER FUNCTION ---
def render_metu_header():
    st.markdown(f"<div class='header-text'>{t['title']}</div>", unsafe_allow_html=True)
    st.markdown(f"""
        <div class='metu-topbar'>{t['topbar']}</div>
        <div class='metu-navbar'>{t['subtitle']}</div>
    """, unsafe_allow_html=True)

# --- LOCAL KNOWLEDGE BASE LOADER ---
@st.cache_data 
def load_knowledge_base():
    kb_text = ""
    kb_folder = "knowledge_base"

    if not os.path.exists(kb_folder):
        return "Warning: 'knowledge_base' folder not found. Please create it and add your files."

    for filename in os.listdir(kb_folder):
        filepath = os.path.join(kb_folder, filename)
        try:
            if filename.endswith(".txt"):
                with open(filepath, "r", encoding="utf-8") as f:
                    kb_text += f"\n\n--- Source: {filename} ---\n" + f.read()
            elif filename.endswith(".pdf"):
                reader = PdfReader(filepath)
                kb_text += f"\n\n--- Source: {filename} ---\n"
                for page in reader.pages:
                    text = page.extract_text()
                    if text: kb_text += text + "\n"
            elif filename.endswith(".pptx"):
                prs = Presentation(filepath)
                kb_text += f"\n\n--- Source: {filename} ---\n"
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            kb_text += shape.text + "\n"
            elif filename.endswith(".docx"):
                doc = Document(filepath)
                kb_text += f"\n\n--- Source: {filename} ---\n"
                for para in doc.paragraphs:
                    if para.text.strip():
                        kb_text += para.text + "\n"
        except Exception as e:
            print(f"Error reading {filename}: {e}")

    return kb_text

# --- INITIALIZE OPENAI SECURELY ---
try:
    client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])
except KeyError:
    st.error("API Key not found! Please configure your Streamlit Secrets.")
    st.stop()

# --- SESSION STATE INITIALIZATION ---
if "messages" not in st.session_state:
    st.session_state.messages = []
    
if "quick_prompt" not in st.session_state:
    st.session_state.quick_prompt = None

# --- SYSTEM INSTRUCTIONS ---
kb_text = load_knowledge_base()

system_instruction = f"""
You are an intelligent virtual consultant specifically designed to assist 
Middle East Technical University (METU) Industrial Engineering students 
with their Summer Practices (IE 300 and IE 400). 

Below is the official documentation. You MUST use ONLY this information to answer user queries. 
DO NOT use outside knowledge. DO NOT invent constraints, deadlines, or rules that are not explicitly written below.

CRITICAL GUARDRAILS:
1. IE 400 COMPANY TYPES: According to the official presentation slides, students can search for alternative companies for IE 400 in ALL of the following sectors:
   - Manufacturing firms (Automotive, Machine parts, Electronics, etc.)
   - Batch process industries (Steel, Paper mills, Pharmaceutical, etc.)
   - Continuous process industries (Cement, concrete, Sugar, Flour mill, etc.)
   - Service industries (Hospitals, Hotels, Banking, Research organizations, Transportation, Public institutions, NGOs, etc.)
   NEVER state that continuous process or service industries are banned for IE 400. They are explicitly accepted.
   For the cases similar to this, do not try to infer information from another SP, i.e. if IE300 is asked just use documents related with IE300, if IE400 is asked just use documents related with IE400 do not mix them.
2. For IE 400 related questions generally use IE400_introduction_2025.pptx and for IE 300 related questions generally use IE300_introduction_2025.pptx.
3. ALWAYS respond in the SAME LANGUAGE as the user's prompt (e.g., reply in Turkish if the prompt is in Turkish).

<official_documentation>
{kb_text}
</official_documentation>
"""

# --- RENDER MAIN UI ---
render_metu_header()

# --- CONTINUE ENHANCED SIDEBAR ---
with st.sidebar:
    st.markdown("---")
    
    # Quick Links Section
    st.header(t["links_header"])
    st.markdown(f"""
        <a href="https://ie.metu.edu.tr/en/current-semester" target="_blank" class="sidebar-link">📚 {t['link_curr']}</a>
        <a href="https://sp-ie.metu.edu.tr/en" target="_blank" class="sidebar-link">🌐 {t['link_sp']}</a>
        <a href="https://ie.metu.edu.tr/en" target="_blank" class="sidebar-link">🏠 {t['link_main']}</a>
    """, unsafe_allow_html=True)
    
    st.markdown("---")
    
    # FAQ Section
    st.header(t["faq_header"])
    
    st.subheader(t["cat1"])
    if st.button(t["q1_lbl"]): st.session_state.quick_prompt = t["q1_prmpt"]
    if st.button(t["q2_lbl"]): st.session_state.quick_prompt = t["q2_prmpt"]
    if st.button(t["q3_lbl"]): st.session_state.quick_prompt = t["q3_prmpt"]
        
    st.subheader(t["cat2"])
    if st.button(t["q4_lbl"]): st.session_state.quick_prompt = t["q4_prmpt"]
    if st.button(t["q5_lbl"]): st.session_state.quick_prompt = t["q5_prmpt"]
    if st.button(t["q6_lbl"]): st.session_state.quick_prompt = t["q6_prmpt"]
    if st.button(t["q7_lbl"]): st.session_state.quick_prompt = t["q7_prmpt"]
    if st.button(t["q8_lbl"]): st.session_state.quick_prompt = t["q8_prmpt"]

    st.markdown("---")
    
    # Contact Information at the bottom
    st.header(t["contact_hdr"])
    st.markdown("""
        **Tel:** +90 (312) 210 4796  
        **Fax:** +90 (312) 210 4786  
        **E-Mail:** [ie-staj@metu.edu.tr](mailto:ie-staj@metu.edu.tr)  
        **Address:** METU Industrial Engineering, 06800
    """)

# --- CHAT RENDERING & LOGIC ---

# 1. Render existing messages
for msg in st.session_state.messages:
    st.chat_message(msg["role"]).write(msg["content"])

# 2. Render welcome recommendations ONLY if chat is empty
if len(st.session_state.messages) == 0:
    st.info(t["welcome_msg"])
    
    # Create a clean 3-column grid for the 6 starter questions
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.markdown('<div class="starter-btn">', unsafe_allow_html=True)
        if st.button(t["rec_1"], use_container_width=True): st.session_state.quick_prompt = t["rec_1_prmpt"]
        if st.button(t["rec_4"], use_container_width=True): st.session_state.quick_prompt = t["rec_4_prmpt"]
        st.markdown('</div>', unsafe_allow_html=True)
        
    with col2:
        st.markdown('<div class="starter-btn">', unsafe_allow_html=True)
        if st.button(t["rec_2"], use_container_width=True): st.session_state.quick_prompt = t["rec_2_prmpt"]
        if st.button(t["rec_5"], use_container_width=True): st.session_state.quick_prompt = t["rec_5_prmpt"]
        st.markdown('</div>', unsafe_allow_html=True)
        
    with col3:
        st.markdown('<div class="starter-btn">', unsafe_allow_html=True)
        if st.button(t["rec_3"], use_container_width=True): st.session_state.quick_prompt = t["rec_3_prmpt"]
        if st.button(t["rec_6"], use_container_width=True): st.session_state.quick_prompt = t["rec_6_prmpt"]
        st.markdown('</div>', unsafe_allow_html=True)

# 3. Handle input
prompt = st.chat_input(t["placeholder"])

if st.session_state.quick_prompt:
    prompt = st.session_state.quick_prompt
    st.session_state.quick_prompt = None

if prompt:
    st.session_state.messages.append({"role": "user", "content": prompt})
    st.chat_message("user").write(prompt)

    # Wrap the API call in a spinner to show the thinking animation
    with st.spinner(t["thinking"]):
        try:
            api_messages = [{"role": "system", "content": system_instruction}] + st.session_state.messages

            response = client.chat.completions.create(
                model="gpt-4o-mini",
                messages=api_messages,
                temperature=0.2
            )
            
            bot_reply = response.choices[0].message.content
            
            st.session_state.messages.append({"role": "assistant", "content": bot_reply})
            st.chat_message("assistant").write(bot_reply)
            
            # Since we appended a new message, rerun to hide the starter recommendations
            st.rerun()
            
        except Exception as e:
            st.error(f"An error occurred: {e}")
